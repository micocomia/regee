import os
import PyPDF2
from pptx import Presentation
from typing import List, Dict, Any, Optional
import logging
from sentence_transformers import SentenceTransformer
import streamlit as st

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """
    Processes documents (PDF, PPTX) for the review chatbot.
    """
    def __init__(self, embedding_model: str = "all-MiniLM-L6-v2", 
                 chunk_size: int = 500, chunk_overlap: int = 100):
        """
        Initialize the document processor.
        
        Args:
            embedding_model: Name of the sentence-transformers model to use
            chunk_size: Size of document chunks in characters
            chunk_overlap: Overlap between chunks in characters
        """
        self.embedding_model = SentenceTransformer(embedding_model)
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
    
    def process_document(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Process a document file and return chunks with embeddings.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            List of document chunks with content, embeddings, and metadata
        """
        # Extract text based on file type
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            page_texts = self._extract_pdf_text(file_path)
            # Extract topics from the combined text of all pages
            full_text = ' '.join([page_info['text'] for page_info in page_texts])
            topics = self._extract_topics(full_text)
            
            # Process each page and track page numbers
            processed_chunks = []
            chunk_id = 0
            
            for page_info in page_texts:
                page_text = page_info['text']
                page_num = page_info['page_number']
                
                # Chunk the page text
                chunks = self._chunk_text(page_text)
                
                # Create embeddings for each chunk from this page
                for chunk in chunks:
                    embedding = self.embedding_model.encode(chunk, show_progress_bar=False)
                    
                    processed_chunks.append({
                        'content': chunk,
                        'embedding': embedding,
                        'metadata': {
                            'source': os.path.basename(file_path),
                            'chunk_id': chunk_id,
                            'topics': topics,
                            'page_number': page_num
                        }
                    })
                    chunk_id += 1
            
            return processed_chunks
            
        elif file_extension in ['.pptx', '.ppt']:
            slide_texts = self._extract_pptx_text(file_path)
            # Extract topics from the combined text of all slides
            full_text = ' '.join([slide_info['text'] for slide_info in slide_texts])
            topics = self._extract_topics(full_text)
            
            # Process each slide and track slide numbers
            processed_chunks = []
            chunk_id = 0
            
            for slide_info in slide_texts:
                slide_text = slide_info['text']
                slide_num = slide_info['slide_number']
                
                # Chunk the slide text
                chunks = self._chunk_text(slide_text)
                
                # Create embeddings for each chunk from this slide
                for chunk in chunks:
                    embedding = self.embedding_model.encode(chunk, show_progress_bar=False)
                    
                    processed_chunks.append({
                        'content': chunk,
                        'embedding': embedding,
                        'metadata': {
                            'source': os.path.basename(file_path),
                            'chunk_id': chunk_id,
                            'topics': topics,
                            'page_number': slide_num  # For slides, use slide number as page number
                        }
                    })
                    chunk_id += 1
            
            return processed_chunks
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")
    
    def _extract_pdf_text(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from PDF file with page tracking.
        
        Returns:
            List of dictionaries with text content and page number
        """
        page_texts = []
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                if text.strip():  # Only add non-empty pages
                    page_texts.append({
                        'text': text,
                        'page_number': page_num + 1  
                    })
        return page_texts
    
    def _extract_pptx_text(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text from PowerPoint file with slide tracking.
        
        Returns:
            List of dictionaries with text content and slide number
        """
        slide_texts = []
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            
            if text.strip():  # Only add non-empty slides
                slide_texts.append({
                    'text': text,
                    'slide_number': slide_num + 1  
                })
        return slide_texts
    
    def _chunk_text(self, text: str) -> List[str]:
        """Split text into chunks with overlap."""
        chunks = []
        start = 0
        
        while start < len(text):
            end = min(start + self.chunk_size, len(text))
            
            # Adjust end to avoid splitting words
            if end < len(text):
                # Look for the last space within the chunk
                last_space = text.rfind(' ', start, end)
                if last_space != -1 and last_space > start:
                   end = last_space + 1  # Include the space

            # Add the chunk
            chunks.append(text[start:end])
            
            # Move the start position for the next chunk, considering overlap
            start = max(end - self.chunk_overlap, start + 1)
        
        return chunks
    
    def _extract_topics(self, text: str) -> List[str]:
        """
        Extract key topics from text (simple version).
        A more sophisticated topic extraction would be implemented here.
        """
        # Simple keyword extraction - in a real system, use TF-IDF or LDA
        common_words = ['the', 'and', 'or', 'to', 'a', 'in', 'that', 'it', 'with']
        words = [word.lower() for word in text.split() if len(word) > 4]
        word_counts = {}
        
        for word in words:
            if word not in common_words:
                word_counts[word] = word_counts.get(word, 0) + 1
        
        # Get the top 5 words as "topics"
        sorted_words = sorted(word_counts.items(), key=lambda x: x[1], reverse=True)
        topics = [word for word, count in sorted_words[:5]]
        
        return topics

# Integration code for Streamlit app
def process_uploaded_file(uploaded_file):
    """Process an uploaded document."""
    # Save the file temporarily
    file_path = os.path.join("./uploads", uploaded_file.name)
    os.makedirs("./uploads", exist_ok=True)
    
    with open(file_path, "wb") as f:
        f.write(uploaded_file.getbuffer())
    
    # Process the document
    document_chunks = document_processor.process_document(file_path)
    
    # Add to vector store
    vector_store.add_documents(document_chunks)
    
    # Extract topics for the UI
    all_topics = set()
    for chunk in document_chunks:
        all_topics.update(chunk['metadata']['topics'])
    
    # Update session state
    st.session_state.documents.append(file_path)
    st.session_state.document_names.append(uploaded_file.name)
    
    for topic in all_topics:
        if topic not in st.session_state.topics:
            st.session_state.topics.append(topic)
            
    return True

# Test the processor
processor = DocumentProcessor()

# Test with a PDF file
pdf_path = r"C:\Users\karel\Downloads\answer-generation-for-retrieval-based-question-answering-systems.pdf" 
# Test if files can be opened
try:
    with open(pdf_path, 'rb') as f:
        print("Successfully opened PDF file")
except Exception as e:
    print(f"Failed to open PDF: {str(e)}")

try:
    from pptx import Presentation
    prs = Presentation(pptx_path)
    print(f"Successfully opened PPTX file with {len(prs.slides)} slides")
except Exception as e:
    print(f"Failed to open PPTX: {str(e)}")

# %%
# Test the processor
print("Initializing DocumentProcessor...")
processor = DocumentProcessor()

# Test with PDF file
pdf_path = r"C:\Users\karel\Downloads\answer-generation-for-retrieval-based-question-answering-systems.pdf"
print(f"Testing PDF processing: {pdf_path}")
try:
    pdf_chunks = processor.process_document(pdf_path)
    print(f"Successfully processed PDF into {len(pdf_chunks)} chunks")
    if pdf_chunks:
        print(f"First chunk content: {pdf_chunks[0]['content'][:100]}...")
        print(f"Topics identified: {pdf_chunks[0]['metadata']['topics']}")
        print(f"Page number: {pdf_chunks[0]['metadata']['page_number']}")
except Exception as e:
    print(f"Error processing PDF: {str(e)}")
    import traceback
    traceback.print_exc()

# Test with PowerPoint file
pptx_path = r"C:\Users\karel\Downloads\DTI 5125 Question Answering Group 1.pptx"
print(f"\nTesting PowerPoint processing: {pptx_path}")
try:
    pptx_chunks = processor.process_document(pptx_path)
    print(f"Successfully processed PowerPoint into {len(pptx_chunks)} chunks")
    if pptx_chunks:
        print(f"First chunk content: {pptx_chunks[0]['content'][:100]}...")
        print(f"Topics identified: {pptx_chunks[0]['metadata']['topics']}")
        print(f"Slide number: {pptx_chunks[0]['metadata']['page_number']}")
except Exception as e:
    print(f"Error processing PPTX: {str(e)}")
    import traceback
    traceback.print_exc()


